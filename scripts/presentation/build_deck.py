"""Build the group-meeting deck from ``content.py``.

    MPLCONFIGDIR=/tmp/matplotlib PYTHONPATH=scripts/presentation \\
        .venv/bin/python scripts/presentation/build_deck.py

Writes two renderings of the same slide model:

* ``Papers/thermal_states_presentation.pptx`` — the editable deliverable.
  Real text runs, real presenter notes, theme fonts the user's PowerPoint has.
* ``Papers/thermal_states_presentation.html`` — a pixel-mirror for review, and
  the route to a PDF (``google-chrome --headless --print-to-pdf``).

Both consume the same absolutely-positioned block list, so what you check in
the HTML is what the pptx contains.
"""
from __future__ import annotations

import html as _html
import math
import re
from pathlib import Path

import equations as EQS
import style as S
from inline_math import Run, parse

REPO = Path(__file__).resolve().parents[2]

# Which deck to build.  Set by main() from the command line; every deck ships a
# content module, its own figure directory, and its own output basename.
DECKS = {
    "results": dict(module="content", figdir="deck",
                    out="thermal_states_presentation",
                    footer="Thermal states as training data for a quantum neuron",
                    lint=False),
    "theory": dict(module="content_theory", figdir="deck_theory",
                   out="molecular_hamiltonians_and_thermal_states",
                   footer="Molecular Hamiltonians and thermal states",
                   lint=True),
    "gap": dict(module="content_gap", figdir="deck_gap",
                out="homo_lumo_gap_diagnosis",
                footer="Why the HOMO–LUMO gap shows no quantum advantage",
                lint=True),
}
DECK = DECKS["results"]
content = None                       # bound by select_deck()
FIGDIR = REPO / "figures" / "deck"

EMU = 914400


def select_deck(name: str) -> None:
    global DECK, content, FIGDIR
    import importlib
    DECK = DECKS[name]
    content = importlib.import_module(DECK["module"])
    FIGDIR = REPO / "figures" / DECK["figdir"]


# --------------------------------------------------------------- block model
def rect(x, y, w, h, fill, line=None, radius=True):
    return dict(t="rect", x=x, y=y, w=w, h=h, fill=fill, line=line, radius=radius)


def text(x, y, w, h, body, size=S.SZ_BODY, color=S.INK, align="l", valign="t",
         font=S.SANS, leading=1.30, space_after=0.0, bullet=False):
    """`body` is a string or a list of strings (one paragraph each)."""
    paras = [body] if isinstance(body, str) else list(body)
    return dict(t="text", x=x, y=y, w=w, h=h, paras=paras, size=size, color=color,
                align=align, valign=valign, font=font, leading=leading,
                space_after=space_after, bullet=bullet)


def img(path, x, y, w, h):
    return dict(t="img", x=x, y=y, w=w, h=h, path=str(path))


# ------------------------------------------------------------ text measuring
# Average advance width as a fraction of the point size, for the body faces.
_ADV = {S.SANS: 0.475, S.SERIF: 0.470}


def wrapped_height(body, w_in, size, leading=1.30, font=S.SANS,
                   space_after=0.0, bullet=False) -> float:
    paras = [body] if isinstance(body, str) else list(body)
    usable = w_in - (0.20 if bullet else 0.0)
    per_line = max(8, int(usable * 72 / (size * _ADV.get(font, 0.475))))
    lines = 0
    for p in paras:
        n = len(_plain(p))
        lines += max(1, math.ceil(n / per_line))
    return lines * size * leading / 72 + max(0, len(paras) - 1) * space_after


def _plain(s: str) -> str:
    return "".join(r.text for r in parse(s))


# ------------------------------------------------------------ shared chrome
def chrome(sl, index, total, dark=False):
    b = []
    if sl.get("kicker"):
        b.append(text(S.MX, S.KICKER_Y, S.CW, 0.24, sl["kicker"].upper(),
                      size=S.KICKER_SZ, color=S.ONDARK if dark else S.GRAY,
                      font=S.SANS, leading=1.0))
    if sl.get("title"):
        b.append(text(S.MX, S.TITLE_Y, S.CW, 0.62, f"**{sl['title']}**",
                      size=S.TITLE_SZ, color=S.WHITE if dark else S.NAVY,
                      font=S.SERIF, leading=1.05))
    if not dark:
        b.append(text(S.MX, S.FOOTER_Y, S.CW * 0.75, 0.22, DECK["footer"],
                      size=S.FOOTER_SZ, color=S.GRAY, leading=1.0))
        b.append(text(S.MX + S.CW * 0.75, S.FOOTER_Y, S.CW * 0.25, 0.22,
                      f"{index} / {total}", size=S.FOOTER_SZ, color=S.GRAY,
                      align="r", leading=1.0))
    return b


def eq_block(name, caption, x, y, w, pt=17.0, fill=S.PANEL):
    """A pale panel holding one display equation, with its explanation under it."""
    path, ew, eh = EQS.render(EQS.EQ[name], pt=pt)
    scale = min(1.0, (w - 0.7) / ew)
    ew, eh = ew * scale, eh * scale
    pad = 0.19
    ph = eh + 2 * pad
    b = [rect(x, y, w, ph, fill),
         img(path, x + (w - ew) / 2, y + pad, ew, eh)]
    yy = y + ph
    if caption:
        ch = wrapped_height(caption, w - 0.08, S.SZ_CAPTION, 1.32)
        b.append(text(x + 0.04, yy + 0.07, w - 0.08, ch + 0.05, caption,
                      size=S.SZ_CAPTION, color=S.SLATE, leading=1.32))
        yy += ch + 0.16
    return b, yy - y


def bullet_block(items, x, y, w, size=S.SZ_BODY, color=S.INK, gap=0.13):
    b, yy = [], y
    for it in items:
        h = wrapped_height(it, w - 0.22, size, 1.34, bullet=True)
        b.append(text(x, yy, 0.16, size / 60, "•", size=size, color=S.BLUE,
                      leading=1.34))
        b.append(text(x + 0.22, yy, w - 0.22, h + 0.04, it, size=size,
                      color=color, leading=1.34))
        yy += h + gap
    return b, yy - y


def stat_row(stats, y, h=0.78):
    b = []
    n = len(stats)
    w = S.CW / n
    for i, (big, small) in enumerate(stats):
        x = S.MX + i * w
        b.append(text(x, y, w, 0.42, f"**{big}**", size=S.SZ_STAT, color=S.NAVY,
                      align="c", font=S.SERIF, leading=1.02))
        b.append(text(x + 0.1, y + 0.46, w - 0.2, 0.32, small, size=S.SZ_SMALL,
                      color=S.SLATE, align="c", leading=1.2))
    return b, h


def note_block(body, x, y, w, fill=S.PANEL2, size=S.SZ_SMALL):
    h = wrapped_height(body, w - 0.5, size, 1.36) + 0.34
    return [rect(x, y, w, h, fill),
            text(x + 0.25, y + 0.15, w - 0.5, h - 0.3, body, size=size,
                 color=S.INK, leading=1.36)], h


def card_row(cards, y, h, x=S.MX, w=S.CW, gap=0.22, dark=False,
             title_color=None, body_color=None):
    b = []
    cw = (w - gap * (len(cards) - 1)) / len(cards)
    for i, (t, body) in enumerate(cards):
        cx = x + i * (cw + gap)
        b.append(rect(cx, y, cw, h, S.NAVY_D if dark else S.PANEL))
        b.append(text(cx + 0.22, y + 0.18, cw - 0.44, 0.30, f"**{t}**",
                      size=S.SZ_CARD_TITLE, font=S.SERIF,
                      color=title_color or (S.WHITE if dark else S.NAVY),
                      leading=1.15))
        items = body if isinstance(body, list) else [body]
        if isinstance(body, list):
            bb, _ = bullet_block(items, cx + 0.22, y + 0.60, cw - 0.44,
                                 size=S.SZ_SMALL,
                                 color=body_color or (S.ONDARK if dark else S.INK))
            b += bb
        else:
            bh = wrapped_height(body, cw - 0.44, S.SZ_SMALL, 1.36)
            b.append(text(cx + 0.22, y + 0.60, cw - 0.44, bh + 0.1, body,
                          size=S.SZ_SMALL, leading=1.36,
                          color=body_color or (S.ONDARK if dark else S.INK)))
    return b


def figure_block(name, x, y, w, caption=None, max_h=None, cap_w=None,
                 cap_x=None, cap_align="l"):
    from PIL import Image
    p = FIGDIR / name
    with Image.open(p) as im:
        ar = im.height / im.width
    h = w * ar
    if max_h and h > max_h:
        h, w = max_h, max_h / ar
    b = [img(p, x, y, w, h)]
    yy = y + h
    if caption:
        cw = cap_w or w
        cx = x if cap_x is None else cap_x
        ch = wrapped_height(caption, cw - 0.1, S.SZ_CAPTION, 1.30)
        b.append(text(cx, yy + 0.09, cw, ch + 0.06, caption, size=S.SZ_CAPTION,
                      color=S.SLATE, align=cap_align, leading=1.30))
        yy += ch + 0.15
    return b, yy - y, w


# ----------------------------------------------------------------- layouts
def L_title(sl, i, n):
    b = [rect(0, 0, S.SLIDE_W, S.SLIDE_H, S.NAVY, radius=False)]
    b.append(text(1.0, 2.05, 10.2, 1.5, f"**{sl['title']}**", size=38,
                  color=S.WHITE, font=S.SERIF, leading=1.14))
    b.append(text(1.0, 3.62, 10.0, 0.8, sl["subtitle"], size=15.5,
                  color=S.ONDARK, leading=1.34))
    b.append(text(1.0, 4.62, 8.0, 0.3, sl["venue"], size=12, color="#8FA0C8",
                  leading=1.0))
    # a half-filled sixteen-wire register
    for k in range(16):
        c = S.CYAN if k % 2 == 0 else "#3E4C86"
        b.append(rect(1.0 + k * 0.40, 5.72, 0.185, 0.185, c, radius="oval"))
    b.append(text(1.0, 6.12, 9.5, 0.28, f"*{sl['strapline']}*", size=10.5,
                  color="#8FA0C8", leading=1.2))
    return b


def L_eq_bullets(sl, i, n):
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    has_fig = bool(sl.get("figure"))
    wide = sl.get("wide_figure", False)
    eq_w = S.CW if (wide or not has_fig) else S.CW
    for name, cap in sl.get("equations", []):
        eb, dh = eq_block(name, cap, S.MX, y, eq_w)
        b += eb
        y += dh + 0.20

    if has_fig and not wide:
        fw = 5.35
        col_w = S.CW - fw - 0.42
        bb, bh = bullet_block(sl["bullets"], S.MX, y + 0.02, col_w)
        b += bb
        cap = sl.get("figure_caption")
        cap_h = (wrapped_height(cap, fw - 0.1, S.SZ_CAPTION, 1.30) + 0.24) if cap else 0
        room = S.FOOTER_Y - 0.24 - y - cap_h - (0.92 if sl.get("stats") else 0)
        fb, fh, fw2 = figure_block(sl["figure"], S.MX + col_w + 0.42, y, fw, cap,
                                   max_h=room)
        for blk in fb:                        # keep the column right-aligned
            blk["x"] += fw - fw2
        b += fb
        y += max(bh, fh) + 0.20
    elif has_fig and wide:
        bb, bh = bullet_block(sl["bullets"], S.MX, y + 0.02, S.CW)
        b += bb
        y += bh + 0.20
        cap = sl.get("figure_caption")
        cap_h = wrapped_height(cap, S.CW - 0.1, S.SZ_CAPTION, 1.30) + 0.24 if cap else 0
        avail = S.FOOTER_Y - 0.22 - y - (0.92 if sl.get("stats") else 0) - cap_h
        fb, fh, fw = figure_block(sl["figure"], S.MX, y, 8.9, None,
                                  max_h=max(1.2, avail))
        for blk in fb:
            blk["x"] += (S.CW - fw) / 2
        b += fb
        y += fh
        if cap:
            b.append(text(S.MX, y + 0.10, S.CW, cap_h, cap, size=S.SZ_CAPTION,
                          color=S.SLATE, align="c", leading=1.30))
            y += cap_h
    else:
        bb, bh = bullet_block(sl["bullets"], S.MX, y + 0.02, S.CW)
        b += bb
        y += bh + 0.18

    if sl.get("note"):
        nb, nh = note_block(sl["note"], S.MX, y, S.CW)
        b += nb
        y += nh + 0.16
    if sl.get("stats"):
        sb, _ = stat_row(sl["stats"], max(y, S.FOOTER_Y - 1.02))
        b += sb
    return b


def L_eq_stack(sl, i, n):
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    for name, cap in sl["equations"]:
        eb, dh = eq_block(name, cap, S.MX, y, S.CW, pt=17.5)
        b += eb
        y += dh + 0.26
    nh = (wrapped_height(sl["note"], S.CW - 0.5, S.SZ_SMALL, 1.36) + 0.34
          if sl.get("note") else 0.0)
    if sl.get("cards"):
        h = min(S.FOOTER_Y - 0.30 - y - (nh + 0.30 if nh else 0.0), 3.1)
        b += card_row(sl["cards"], y, h)
        y += h + 0.30
    if nh:
        nb, _ = note_block(sl["note"], S.MX, y, S.CW)
        b += nb
    return b


def L_dictionary(sl, i, n):
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    left_w = 7.15
    row_h = 0.53
    for k, (term, meaning) in enumerate(sl["rows"]):
        yy = y + k * (row_h + 0.10)
        b.append(rect(S.MX, yy, 2.95, row_h, S.PANEL))
        b.append(text(S.MX + 0.18, yy + 0.13, 2.6, row_h - 0.2, f"**{term}**",
                      size=11.0, color=S.NAVY, leading=1.16))
        b.append(text(S.MX + 3.20, yy + 0.14, left_w - 3.20, row_h - 0.2,
                      meaning, size=11.0, color=S.INK, leading=1.16))
    lh = len(sl["rows"]) * (row_h + 0.10)

    px = S.MX + left_w + 0.35
    pw = S.CW - left_w - 0.35
    b.append(rect(px, y, pw, lh - 0.10, S.NAVY))
    b.append(text(px + 0.26, y + 0.20, pw - 0.52, 0.28,
                  f"**{sl['panel_title']}**", size=13.0, font=S.SERIF,
                  color=S.WHITE, leading=1.15))
    pb, ph = bullet_block(sl["panel_bullets"], px + 0.26, y + 0.62, pw - 0.52,
                          size=S.SZ_SMALL, color=S.ONDARK, gap=0.08)
    for blk in pb:                       # bullet dots on dark ground
        if blk["t"] == "text" and blk["paras"] == ["•"]:
            blk["color"] = S.CYAN
    b += pb
    nh = wrapped_height(sl["panel_note"], pw - 0.52, 10.0, 1.34)
    b.append(text(px + 0.26, y + 0.62 + ph + 0.10, pw - 0.52, nh + 0.1,
                  sl["panel_note"], size=10.0, color="#9FB3E0", leading=1.34))

    y += lh + 0.16
    eb, _ = eq_block(sl["equation"], sl["equation_caption"], S.MX, y, S.CW,
                     pt=16.5)
    b += eb
    return b


def L_three_faces(sl, i, n):
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    cw = 3.05
    cx = S.MX
    t, eq, sub = sl["centre"]
    nh = (wrapped_height(sl["note"], S.CW - 0.5, S.SZ_SMALL, 1.36) + 0.34
          if sl.get("note") else 0.0)
    CARD_H = min(3.9, S.FOOTER_Y - 0.30 - y - (nh + 0.30 if nh else 0.0))
    b.append(rect(cx, y, cw, CARD_H, S.NAVY))
    b.append(text(cx + 0.24, y + 0.30, cw - 0.48, 0.32, f"**{t}**", size=15.0,
                  font=S.SERIF, color=S.WHITE, align="c", leading=1.1))
    path, ew, eh = EQS.render(eq.strip("$"), pt=22.0, color=S.CYAN)
    b.append(img(path, cx + (cw - ew) / 2, y + 0.86, ew, eh))
    b.append(text(cx + 0.24, y + 1.62, cw - 0.48, 0.5, sub, size=10.5,
                  color=S.ONDARK, align="c", leading=1.28))
    b.append(text(cx + 0.24, y + CARD_H - 0.52, cw - 0.48, 0.3,
                  "stored once, per state", size=9.5, color="#8FA0C8",
                  align="c", leading=1.1))

    fx = S.MX + cw + 0.46
    fw = S.CW - cw - 0.46
    each = (fw - 2 * 0.20) / 3
    for k, (title, body, tag) in enumerate(sl["faces"]):
        x = fx + k * (each + 0.20)
        b.append(rect(x, y, each, CARD_H, S.PANEL))
        b.append(text(x + 0.20, y + 0.18, each - 0.40, 0.32, f"**{title}**",
                      size=12.0, font=S.SERIF, color=S.NAVY, leading=1.14))
        bh = wrapped_height(body, each - 0.40, 10.2, 1.34)
        b.append(text(x + 0.20, y + 0.62, each - 0.40, bh + 0.1, body,
                      size=10.2, color=S.INK, leading=1.34))
        b.append(text(x + 0.20, y + CARD_H - 0.46, each - 0.40, 0.26, f"*{tag}*",
                      size=9.5, color=S.RUST, leading=1.1))
    # connectors
    for k in range(3):
        x = fx + k * (each + 0.20)
        b.append(rect(S.MX + cw + 0.10, y + CARD_H / 2, 0.26, 0.028, S.GRAY,
                      radius=False))
        break
    y += CARD_H + 0.30
    nb, nh = note_block(sl["note"], S.MX, y, S.CW)
    b += nb
    return b


def L_steps_eq(sl, i, n):
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    for name, cap in sl["equations"]:
        eb, dh = eq_block(name, cap, S.MX, y, S.CW, pt=17.0)
        b += eb
        y += dh + 0.20
    has_fig = bool(sl.get("figure"))
    col_w = (S.CW - 5.4 - 0.4) if has_fig else S.CW

    # The note is bottom-anchored; the steps take the space that is left, with
    # their spacing compressed rather than allowed to run into the footer.
    nh = (wrapped_height(sl["note"], S.CW - 0.5, S.SZ_SMALL, 1.36) + 0.34
          if sl.get("note") else 0.0)
    floor = S.FOOTER_Y - 0.16 - (nh + 0.16 if nh else 0.0)
    bodies = [wrapped_height(body, col_w - 0.44, S.SZ_SMALL, 1.34)
              for _, body in sl["steps"]]
    fixed = sum(0.30 + h for h in bodies)
    gap = max(0.08, min(0.22, (floor - y - fixed) / max(1, len(bodies))))

    yy = y
    for k, ((t, body), bh) in enumerate(zip(sl["steps"], bodies)):
        b.append(rect(S.MX, yy + 0.02, 0.30, 0.30, S.NAVY))
        b.append(text(S.MX, yy + 0.055, 0.30, 0.26, f"**{k+1}**", size=11.0,
                      color=S.WHITE, align="c", leading=1.1))
        b.append(text(S.MX + 0.44, yy, col_w - 0.44, 0.28, f"**{t}**",
                      size=12.5, font=S.SERIF, color=S.NAVY, leading=1.15))
        b.append(text(S.MX + 0.44, yy + 0.30, col_w - 0.44, bh + 0.1, body,
                      size=S.SZ_SMALL, color=S.INK, leading=1.34))
        yy += 0.30 + bh + gap
    if has_fig:
        fb, _, _ = figure_block(sl["figure"], S.MX + col_w + 0.4, y, 5.4,
                                sl.get("figure_caption"))
        b += fb
    if nh:
        nb, _ = note_block(sl["note"], S.MX, max(yy + 0.06, floor + 0.16), S.CW)
        b += nb
    return b


def L_two_cards_eq(sl, i, n):
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    b += card_row(sl["cards"], y, 1.98)
    y += 1.98 + 0.28
    for name, cap in sl["equations"]:
        eb, dh = eq_block(name, cap, S.MX, y, S.CW, pt=16.5)
        b += eb
        y += dh + 0.22
    nb, nh = note_block(sl["note"], S.MX, y, S.CW)
    b += nb
    return b


def L_figure_hero(sl, i, n):
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    for name, cap in sl.get("equations", []):
        eb, dh = eq_block(name, cap, S.MX, y, S.CW, pt=15.5)
        b += eb
        y += dh + 0.18

    right = sl.get("callout") or sl.get("bullets")
    fw = 7.35 if right else 9.3
    fx = S.MX if right else S.MX + (S.CW - fw) / 2
    bottom = S.FOOTER_Y - 0.30 - (0.92 if sl.get("stats") else 0.0)
    fb, fh, fw2 = figure_block(sl["figure"], fx, y, fw, sl.get("figure_caption"),
                               max_h=bottom - y - 0.5)
    b += fb
    if right:
        rx = S.MX + fw + 0.42
        rw = S.CW - fw - 0.42
        if sl.get("callout"):
            t, body = sl["callout"]
            hh = wrapped_height(body, rw - 0.5, S.SZ_SMALL, 1.36) + \
                wrapped_height(t, rw - 0.5, S.SZ_CARD_TITLE, 1.16) + 0.58
            b.append(rect(rx, y, rw, min(hh, bottom - y), S.PANEL))
            b.append(text(rx + 0.25, y + 0.20, rw - 0.5, 0.34, f"**{t}**",
                          size=S.SZ_CARD_TITLE, font=S.SERIF, color=S.RUST,
                          leading=1.16))
            th = wrapped_height(t, rw - 0.5, S.SZ_CARD_TITLE, 1.16)
            b.append(text(rx + 0.25, y + 0.24 + th, rw - 0.5,
                          bottom - y - 0.4 - th, body, size=S.SZ_SMALL,
                          color=S.INK, leading=1.36))
        else:
            bb, _ = bullet_block(sl["bullets"], rx, y + 0.02, rw,
                                 size=S.SZ_SMALL)
            b += bb
    if sl.get("stats"):
        sb, _ = stat_row(sl["stats"], S.FOOTER_Y - 1.00)
        b += sb
    return b


def L_summary(sl, i, n):
    b = [rect(0, 0, S.SLIDE_W, S.SLIDE_H, S.NAVY, radius=False)]
    b += chrome(sl, i, n, dark=True)
    y = S.BODY_TOP + 0.10
    h = 4.62
    cw = (S.CW - 2 * 0.28) / 3
    for k, (t, items) in enumerate(sl["columns"]):
        x = S.MX + k * (cw + 0.28)
        b.append(rect(x, y, cw, h, "#28336F"))
        b.append(text(x + 0.26, y + 0.24, cw - 0.52, 0.32, f"**{t}**",
                      size=15.0, font=S.SERIF, color=S.WHITE, leading=1.12))
        bb, _ = bullet_block(items, x + 0.26, y + 0.74, cw - 0.52,
                             size=10.6, color=S.ONDARK, gap=0.16)
        for blk in bb:
            if blk["t"] == "text" and blk["paras"] == ["•"]:
                blk["color"] = S.CYAN
        b += bb
    b.append(text(S.MX, y + h + 0.30, S.CW, 0.34, f"*{sl['strapline']}*",
                  size=11.5, color="#9FB3E0", align="c", leading=1.2))
    return b


def L_references(sl, i, n):
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    for title, items in sl["groups"]:
        b.append(text(S.MX, y, 2.3, 0.28, f"**{title}**", size=11.5,
                      font=S.SERIF, color=S.RUST, leading=1.15))
        bb, bh = bullet_block(items, S.MX + 2.45, y - 0.02, S.CW - 2.45,
                              size=10.6, gap=0.07)
        b += bb
        y += max(bh, 0.34) + 0.24
    nb, _ = note_block(sl["note"], S.MX, max(y + 0.06, S.FOOTER_Y - 1.30), S.CW,
                       size=10.4)
    b += nb
    return b


def L_flow(sl, i, n):
    """A left-to-right chain of labelled stages, then supporting text."""
    b = chrome(sl, i, n)
    y = S.BODY_TOP + 0.10
    stages = sl["stages"]
    gap, arrow = 0.30, 0.26
    sw = (S.CW - (len(stages) - 1) * (gap + arrow)) / len(stages)
    h = 1.30
    for k, (t, sub) in enumerate(stages):
        x = S.MX + k * (sw + gap + arrow)
        dark = k in sl.get("highlight", ())
        b.append(rect(x, y, sw, h, S.NAVY if dark else S.PANEL))
        b.append(text(x + 0.12, y + 0.22, sw - 0.24, 0.52, f"**{t}**", size=11.5,
                      font=S.SERIF, align="c", leading=1.16,
                      color=S.WHITE if dark else S.NAVY))
        b.append(text(x + 0.10, y + 0.78, sw - 0.20, 0.42, sub, size=9.4,
                      align="c", leading=1.22,
                      color=S.ONDARK if dark else S.SLATE))
        if k < len(stages) - 1:
            b.append(text(x + sw + gap / 2 - 0.03, y + h / 2 - 0.16, arrow + 0.06,
                          0.32, "→", size=15.0, color=S.GRAY, align="c",
                          leading=1.0))
    y += h + 0.34
    if sl.get("caption"):
        ch = wrapped_height(sl["caption"], S.CW, S.SZ_CAPTION, 1.30)
        b.append(text(S.MX, y, S.CW, ch + 0.06, sl["caption"], size=S.SZ_CAPTION,
                      color=S.SLATE, align="c", leading=1.30))
        y += ch + 0.30
    if sl.get("cards"):
        b += card_row(sl["cards"], y, min(2.35, S.FOOTER_Y - 0.30 - y))
        y += min(2.35, S.FOOTER_Y - 0.30 - y) + 0.20
    if sl.get("bullets"):
        bb, bh = bullet_block(sl["bullets"], S.MX, y, S.CW)
        b += bb
        y += bh + 0.16
    if sl.get("note"):
        nh = wrapped_height(sl["note"], S.CW - 0.5, S.SZ_SMALL, 1.36) + 0.34
        nb, _ = note_block(sl["note"], S.MX, min(y, S.FOOTER_Y - 0.16 - nh), S.CW)
        b += nb
    return b


def L_table(sl, i, n):
    """A header row and body rows, ruled rather than boxed."""
    b = chrome(sl, i, n)
    y = S.BODY_TOP
    if sl.get("lede"):
        lh = wrapped_height(sl["lede"], S.CW, S.SZ_BODY, 1.34)
        b.append(text(S.MX, y, S.CW, lh + 0.06, sl["lede"], size=S.SZ_BODY,
                      color=S.INK, leading=1.34))
        y += lh + 0.30

    nh = (wrapped_height(sl["note"], S.CW - 0.5, S.SZ_SMALL, 1.36) + 0.34
          if sl.get("note") else 0.0)
    room = S.FOOTER_Y - 0.20 - y - 0.40 - (nh + 0.26 if nh else 0.0)

    cols = sl["columns"]                       # [(heading, relative width)]
    tot = sum(w for _, w in cols)
    xs, x = [], S.MX
    for _, w in cols:
        xs.append((x, S.CW * w / tot))
        x += S.CW * w / tot
    b.append(rect(S.MX, y, S.CW, 0.40, S.NAVY))
    for (head, _), (cx, cw) in zip(cols, xs):
        b.append(text(cx + 0.16, y + 0.11, cw - 0.28, 0.24, f"**{head}**",
                      size=10.4, color=S.WHITE, leading=1.1))
    y += 0.40

    def measure(size, pad):
        return [max(wrapped_height(c, cw - 0.32, size, 1.28)
                    for c, (_, cw) in zip(row, xs)) + pad for row in sl["rows"]]

    size, pad = S.SZ_SMALL, 0.26
    while sum(measure(size, pad)) > room and size > 8.4:
        size, pad = size - 0.3, max(0.16, pad - 0.015)
    for r, (row, rh) in enumerate(zip(sl["rows"], measure(size, pad))):
        if r % 2 == 0:
            b.append(rect(S.MX, y, S.CW, rh, S.PANEL, radius=False))
        for c, (cx, cw) in zip(row, xs):
            b.append(text(cx + 0.16, y + pad / 2, cw - 0.32, rh - pad + 0.06, c,
                          size=size, color=S.INK, leading=1.28))
        y += rh
    if nh:
        nb, _ = note_block(sl["note"], S.MX, y + 0.26, S.CW)
        b += nb
    return b


LAYOUTS = {"title": L_title, "eq_bullets": L_eq_bullets, "eq_stack": L_eq_stack,
           "dictionary": L_dictionary, "three_faces": L_three_faces,
           "steps_eq": L_steps_eq, "two_cards_eq": L_two_cards_eq,
           "figure_hero": L_figure_hero, "summary": L_summary,
           "references": L_references, "flow": L_flow, "table": L_table}


# Layouts that stack content from the top and can therefore leave a large gap
# above the footer.  A slide that is well under-full reads better nudged down;
# full-bleed layouts and the figure-led one are left alone.
_BALANCE = {"eq_bullets", "eq_stack", "steps_eq", "two_cards_eq", "three_faces",
            "dictionary", "table", "flow"}


def _balance(blocks):
    """Shift content down by part of the unused height, for optical balance."""
    body = [b for b in blocks
            if S.BODY_TOP - 0.15 <= b["y"] < S.FOOTER_Y - 0.01]
    if not body:
        return blocks
    slack = (S.FOOTER_Y - 0.26) - max(b["y"] + b["h"] for b in body)
    if slack < 0.40:
        return blocks
    shift = min(0.45 * slack, 1.00)
    for b in body:
        b["y"] += shift
    return blocks


def build_blocks():
    n = len(content.SLIDES)
    out = []
    for i, sl in enumerate(content.SLIDES):
        blocks = LAYOUTS[sl["layout"]](sl, i + 1, n)
        if sl["layout"] in _BALANCE:
            blocks = _balance(blocks)
        out.append((sl, blocks))
    return out


# --------------------------------------------------------------- pptx output
def render_pptx(deck) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Pt

    def rgb(h):
        return RGBColor.from_string(h.lstrip("#").upper())

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(int(S.SLIDE_W * EMU)), Emu(int(S.SLIDE_H * EMU))
    blank = prs.slide_layouts[6]
    ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}

    for sl, blocks in deck:
        s = prs.slides.add_slide(blank)
        for b in blocks:
            if b["t"] == "rect":
                from pptx.enum.shapes import MSO_SHAPE
                kind = {"oval": MSO_SHAPE.OVAL,
                        True: MSO_SHAPE.ROUNDED_RECTANGLE,
                        False: MSO_SHAPE.RECTANGLE}[b["radius"]]
                shp = s.shapes.add_shape(
                    kind, Emu(int(b["x"] * EMU)), Emu(int(b["y"] * EMU)),
                    Emu(int(b["w"] * EMU)), Emu(int(b["h"] * EMU)))
                shp.fill.solid()
                shp.fill.fore_color.rgb = rgb(b["fill"])
                shp.line.fill.background()
                shp.shadow.inherit = False
                if b["radius"] is True:
                    # PowerPoint measures the adjustment against half the SHORT
                    # side, so this reproduces the HTML's fixed 0.055 in radius.
                    try:
                        shp.adjustments[0] = min(
                            0.5, S.RADIUS / (min(b["w"], b["h"]) / 2))
                    except (IndexError, KeyError):
                        pass
            elif b["t"] == "img":
                s.shapes.add_picture(b["path"], Emu(int(b["x"] * EMU)),
                                     Emu(int(b["y"] * EMU)),
                                     Emu(int(b["w"] * EMU)), Emu(int(b["h"] * EMU)))
            elif b["t"] == "text":
                tb = s.shapes.add_textbox(Emu(int(b["x"] * EMU)), Emu(int(b["y"] * EMU)),
                                          Emu(int(b["w"] * EMU)), Emu(int(b["h"] * EMU)))
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                tf.vertical_anchor = (MSO_ANCHOR.MIDDLE if b["valign"] == "m"
                                      else MSO_ANCHOR.TOP)
                for k, para in enumerate(b["paras"]):
                    p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                    p.alignment = ALIGN[b["align"]]
                    p.line_spacing = b["leading"]
                    if b["space_after"]:
                        p.space_after = Pt(b["space_after"] * 72)
                    italic_all = para.startswith("*") and para.endswith("*") \
                        and not para.startswith("**")
                    src = para.strip("*") if italic_all else para
                    for r in parse(src):
                        run = p.add_run()
                        run.text = r.text
                        f = run.font
                        f.size = Pt(b["size"] * (0.72 if r.baseline else 1.0))
                        f.name = b["font"]
                        f.bold = r.bold
                        f.italic = r.italic or italic_all
                        f.color.rgb = rgb(b["color"])
                        if r.baseline:
                            f._rPr.set("baseline",
                                       "30000" if r.baseline == "sup" else "-25000")
        if sl.get("notes"):
            s.notes_slide.notes_text_frame.text = sl["notes"].strip()

    out = REPO / "Papers" / f"{DECK['out']}.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


# --------------------------------------------------------------- html mirror
CSS = """
*{box-sizing:border-box;margin:0;padding:0}
body{background:#6A7080;font-family:%s;-webkit-font-smoothing:antialiased}
.slide{position:relative;width:%.4fin;height:%.4fin;background:#fff;
  overflow:hidden;margin:0 auto 0.22in;page-break-after:always;box-shadow:0 2px 14px rgba(0,0,0,.3)}
.b{position:absolute}
.r{border-radius:0.055in}
.o{border-radius:50%%}
.t{white-space:pre-wrap;word-wrap:break-word}
sup{font-size:.72em;vertical-align:super;line-height:0}
sub{font-size:.72em;vertical-align:sub;line-height:0}
img{display:block;width:100%%;height:100%%}
@page{size:%.4fin %.4fin;margin:0}
@media print{body{background:#fff}.slide{margin:0;box-shadow:none}}
""" % (S.SANS_STACK, S.SLIDE_W, S.SLIDE_H, S.SLIDE_W, S.SLIDE_H)


def _runs_html(src, italic_all=False):
    out = []
    for r in parse(src):
        t = _html.escape(r.text)
        if r.baseline == "sup":
            t = f"<sup>{t}</sup>"
        elif r.baseline == "sub":
            t = f"<sub>{t}</sub>"
        if r.italic or italic_all:
            t = f"<i>{t}</i>"
        if r.bold:
            t = f"<b>{t}</b>"
        out.append(t)
    return "".join(out)


def render_html(deck) -> Path:
    parts = [f"<!doctype html><meta charset=utf-8><title>{content.TITLE}</title>",
             f"<style>{CSS}</style>"]
    for sl, blocks in deck:
        parts.append('<section class="slide">')
        for b in blocks:
            st = (f'left:{b["x"]:.4f}in;top:{b["y"]:.4f}in;'
                  f'width:{b["w"]:.4f}in;height:{b["h"]:.4f}in;')
            if b["t"] == "rect":
                cls = {"oval": "b o", True: "b r", False: "b"}[b["radius"]]
                parts.append(f'<div class="{cls}" style="{st}background:{b["fill"]}"></div>')
            elif b["t"] == "img":
                rel = Path(b["path"]).resolve()
                parts.append(f'<div class="b" style="{st}">'
                             f'<img src="file://{rel}" alt=""></div>')
            else:
                fam = S.SERIF_STACK if b["font"] == S.SERIF else S.SANS_STACK
                al = {"l": "left", "c": "center", "r": "right"}[b["align"]]
                fl = ("display:flex;flex-direction:column;justify-content:center;"
                      if b["valign"] == "m" else "")
                inner = []
                for para in b["paras"]:
                    it = para.startswith("*") and para.endswith("*") and \
                        not para.startswith("**")
                    inner.append("<div>" + _runs_html(para.strip("*") if it else para,
                                                      it) + "</div>")
                parts.append(
                    f'<div class="b t" style="{st}font-family:{fam};'
                    f'font-size:{b["size"]}pt;line-height:{b["leading"]};'
                    f'color:{b["color"]};text-align:{al};{fl}">'
                    + "".join(inner) + "</div>")
        parts.append("</section>")
    out = REPO / "Papers" / f"{DECK['out']}.html"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text("\n".join(parts), encoding="utf-8")
    return out


# --------------------------------------------------------------------- lint
# A deck may declare ``lint=True`` to forbid the em dash in body text and in
# speaker notes.  A dash is easy to reintroduce by accident during an edit, so
# the build checks rather than trusting the author.
BANNED = {"—": "em dash", "–": "en dash"}
# An en dash is legitimate inside a hyphenated proper name (Fermi-Dirac,
# Born-Oppenheimer, Jordan-Wigner) and inside numeric ranges.
EN_DASH_OK = re.compile(r"(?<=[A-Za-z0-9])–(?=[A-Za-z0-9])")

# Constructions that read as informal or evasive in an academic register.  The
# slide text should assert what is true; the contrastive forms below say what is
# false first and are almost always replaceable by a direct statement.  These
# are reported for review, not treated as errors, because a few are legitimate.
STYLE = [
    (re.compile(r"\bnot\b[^.;]{0,40}?\bbut\b"), "not X but Y"),
    (re.compile(r"\brather than\b"), "rather than"),
    (re.compile(r"\bis not\b|\bare not\b|\bdoes not\b|\bdo not\b"), "negation"),
    (re.compile(r"\bIt is worth\b|\bworth noting\b|\bin practice\b"), "filler"),
    (re.compile(r"\bturns out\b|\bthe point is\b|\bof course\b"), "informal"),
]


def lint_text(deck, style: bool = True) -> list[str]:
    problems = []

    def check(where: str, s: str) -> None:
        if not isinstance(s, str):
            return
        if "—" in s:
            problems.append(f"{where}: em dash in {s[:70]!r}")
        for m in re.finditer("–", s):
            if not EN_DASH_OK.match(s, m.start()):
                problems.append(f"{where}: bare en dash in {s[:70]!r}")
        if style and ".notes" not in where:      # slide text only
            for pat, name in STYLE:
                m = pat.search(s)
                if m:
                    problems.append(f"{where}: [{name}] {s[max(0,m.start()-30):m.end()+30]!r}")

    def walk(where, obj):
        if isinstance(obj, str):
            check(where, obj)
        elif isinstance(obj, dict):
            for k, v in obj.items():
                walk(f"{where}.{k}", v)
        elif isinstance(obj, (list, tuple)):
            for k, v in enumerate(obj):
                walk(f"{where}[{k}]", v)

    for i, (sl, _) in enumerate(deck, 1):
        walk(f"slide {i}", sl)
    return problems


def main(argv=None) -> None:
    import argparse
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("deck", nargs="?", default="results", choices=sorted(DECKS),
                    help="which presentation to build")
    ap.add_argument("--no-lint", action="store_true",
                    help="skip the punctuation check")
    a = ap.parse_args(argv)

    select_deck(a.deck)
    deck = build_blocks()
    print(f"{a.deck}: {len(deck)} slides")

    if DECK.get("lint") and not a.no_lint:
        for p in lint_text(deck):
            print(f"  ! {p}")

    for i, (sl, blocks) in enumerate(deck, 1):
        # content must clear the footer rule; full-bleed backgrounds and the
        # footer itself are exempt.
        over = [b for b in blocks
                if 0.05 < b["y"] < S.FOOTER_Y - 0.01 and b["x"] > 0.05
                and (b["y"] + b["h"] > S.FOOTER_Y - 0.04
                     or b["x"] + b["w"] > S.SLIDE_W - 0.05)]
        if over:
            print(f"  ! slide {i:2d} '{sl.get('title', sl['layout'])[:44]}': "
                  f"{len(over)} block(s) reach {max(b['y']+b['h'] for b in over):.2f} in "
                  f"(footer at {S.FOOTER_Y})")
    print("pptx ->", render_pptx(deck))
    print("html ->", render_html(deck))


if __name__ == "__main__":
    main()
